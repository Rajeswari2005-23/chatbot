import torch
from transformers import AutoTokenizer, BitsAndBytesConfig, AutoModelForCausalLM
from peft import PeftModel
import gradio as gr
import argparse
import logging
import psycopg2
from sentence_transformers import SentenceTransformer
from pgvector.psycopg2 import register_vector
import numpy as np
import re
import inflect
import json
import difflib
import docx
import PyPDF2
import os
import speech_recognition as sr
import pyttsx3
import tempfile
import wave
import pyaudio
import threading
import time
from gtts import gTTS
import io
import base64
import pandas as pd  # For Excel and CSV
from PIL import Image  # For image OCR
import pytesseract  # For image OCR
import pptx  # For pptx
import textract  # For .doc, .ppt, etc.

document_context = None

# 1. Parse Arguments
parser = argparse.ArgumentParser(description="Run chatbot with pgvector and structured memory")
parser.add_argument("--model_dir", type=str, default="./outputs/final_model", help="Path to fine-tuned model dir")
args = parser.parse_args()

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# 2. Load Tokenizer
logger.info("Loading tokenizer...")
tokenizer = AutoTokenizer.from_pretrained(args.model_dir, trust_remote_code=True)
tokenizer.pad_token = tokenizer.eos_token

# 3. Quantization Config
bnb_config = BitsAndBytesConfig(
    load_in_4bit=True,
    bnb_4bit_compute_dtype=torch.float16,
    bnb_4bit_use_double_quant=True,
    bnb_4bit_quant_type="nf4"
)

# 4. Load Fine-tuned Model
logger.info("Loading fine-tuned model...")
base_model = AutoModelForCausalLM.from_pretrained(
    "teknium/OpenHermes-2.5-Mistral-7B",
    quantization_config=bnb_config,
    device_map="auto",
    trust_remote_code=True
)
finetuned_model = PeftModel.from_pretrained(base_model, args.model_dir)
finetuned_model.eval()

# 5. PostgreSQL Setup
conn = psycopg2.connect(
    host="localhost",
    database="postgres",
    user="postgres",
    password="chatbot"
)
register_vector(conn)
cursor = conn.cursor()
cursor.execute("""
    CREATE EXTENSION IF NOT EXISTS vector;
    CREATE TABLE IF NOT EXISTS memory_vectors (
        id SERIAL PRIMARY KEY,
        content TEXT,
        embedding VECTOR(384)
    );
    CREATE TABLE IF NOT EXISTS relationship_memory (
        id SERIAL PRIMARY KEY,
        content TEXT,
        embedding VECTOR(384)
    );
    CREATE TABLE IF NOT EXISTS document_memory (
        id SERIAL PRIMARY KEY,
        question TEXT,
        answer TEXT
    );
""")
conn.commit()

# 6. Embedding Model
embedding_model = SentenceTransformer("all-MiniLM-L6-v2")

# 7. Audio Setup
recognizer = sr.Recognizer()
engine = pyttsx3.init()

# Configure text-to-speech engine
voices = engine.getProperty('voices')
if voices:
    engine.setProperty('voice', voices[0].id)  # Use first available voice
engine.setProperty('rate', 150)  # Speed of speech
engine.setProperty('volume', 0.9)  # Volume level

# Audio recording parameters
CHUNK = 1024
FORMAT = pyaudio.paInt16
CHANNELS = 1
RATE = 44100
RECORD_SECONDS = 5

# Global variables for audio recording
recording = False
audio_frames = []
p = inflect.engine()  # For pluralization
audio_interface = pyaudio.PyAudio()  # For audio recording

def get_embedding(text):
    return np.array(embedding_model.encode(text)).tolist()


def insert_memory(text):
    embedding = get_embedding(text)
    cursor.execute(
        "INSERT INTO memory_vectors (content, embedding) VALUES (%s, %s)",
        (text, embedding)
    )
    conn.commit()
    print(f"✅ Inserted semantic memory: {text}")


def insert_relationship(text):
    embedding = get_embedding(text)
    cursor.execute(
        "INSERT INTO relationship_memory (content, embedding) VALUES (%s, %s)",
        (text, embedding)
    )
    conn.commit()


def search_similar_memory(query, top_k=3):
    embedding = get_embedding(query)
    cursor.execute(
        "SELECT content FROM memory_vectors ORDER BY embedding <-> %s::vector LIMIT %s",
        (embedding, top_k)
    )
    results = cursor.fetchall()
    return [r[0] for r in results]


def search_relationship_memory(query, top_k=5):
    embedding = get_embedding(query)
    cursor.execute(
        "SELECT content FROM relationship_memory ORDER BY embedding <-> %s::vector LIMIT %s",
        (embedding, top_k)
    )
    results = cursor.fetchall()
    return [r[0] for r in results]


def is_fact_statement(text):
    #return '?' not in text and " is " in text.lower()
    q_keywords=["how","what","where","which","who","may","will","for","would","can","had","has","does","might","could","whose","whom","was","shall","should","were","are","did","do","when","why","have","is","am","in","on"]
    iq_keywords=["list","explain","describe","define","discuss","illustrate","regarding","find","tell","compare","analyse","evaluate","justify","outline","summarize","give","state","mention","choose"]
    List=list(text.split())
    return '?'not in  text and List[0].lower() not in q_keywords and List[0].lower() not in iq_keywords


def normalize_input(user_input):
    return user_input


def fetch_structured_response(user_input):
    patterns = [
        r"(?:Who is|Who are|What is|What are) ([^?]+?)['']s ([^?]+)",
        r"(?:Who is|Who are|What is|What are) the ([^?]+) of ([^?]+)",
        r"(?:Who is|Who are|What is|What are) my ([^?]+)",
    ]
    for pattern in patterns:
        match = re.match(pattern, user_input.strip(), re.IGNORECASE)
        if match:
            if len(match.groups()) == 2:
                if " of " in pattern:
                    relation = match.group(1).strip().lower()
                    subject = match.group(2).strip().capitalize()
                else:
                    subject = match.group(1).strip().capitalize()
                    relation = match.group(2).strip().lower()
            else:
                subject = "Rajeswari"
                relation = match.group(1).strip().lower()

            singular = p.singular_noun(relation)
            plural = p.plural(relation)
            relations_to_try = {relation}
            if singular:
                relations_to_try.add(singular)
            if plural:
                relations_to_try.add(plural)

            print(f"🔍 Trying: subject={subject}, relations={relations_to_try}")
            for rel in relations_to_try:
                cursor.execute(
                    "SELECT object FROM relationship_memory WHERE subject = %s AND relation = %s",
                    (subject, rel)
                )
                results = cursor.fetchall()
                if results:
                    objects = [r[0] for r in results]
                    return f"{subject}'s {rel} is {', '.join(objects)}."
    return None


# Load your dataset (test.json) once at startup
with open("test.json", "r", encoding="utf-8") as f:
    dataset_qa = json.load(f)

def find_dataset_item(user_input):
    user_input_lower = user_input.strip().lower()
    instructions = [item["instruction"].strip().lower() for item in dataset_qa]
    # Exact match
    for idx, instr in enumerate(instructions):
        if instr == user_input_lower:
            return dataset_qa[idx]
    # Substring match
    for idx, instr in enumerate(instructions):
        if instr in user_input_lower or user_input_lower in instr:
            return dataset_qa[idx]
    # Closest fuzzy match
    close_matches = difflib.get_close_matches(user_input_lower, instructions, n=1, cutoff=0.0)
    if close_matches:
        idx = instructions.index(close_matches[0])
        return dataset_qa[idx]
    # Fallback to first item if nothing matches
    return dataset_qa[0]


# Audio Functions
def start_recording():
    global recording, audio_frames
    recording = True
    audio_frames = []
    print("DEBUG: Starting audio recording...")

    def record_audio():
        global recording, audio_frames
        try:
            stream = audio_interface.open(format=FORMAT,
                                         channels=CHANNELS,
                                         rate=RATE,
                                         input=True,
                                         frames_per_buffer=CHUNK)
            print("DEBUG: Audio stream opened.")
            while recording:
                data = stream.read(CHUNK, exception_on_overflow=False)
                audio_frames.append(data)
            stream.stop_stream()
            stream.close()
            print("DEBUG: Audio stream closed.")
        except Exception as e:
            print(f"ERROR: Exception in record_audio: {e}")

    thread = threading.Thread(target=record_audio)
    thread.start()
    return "🎤 Recording started... Click 'Stop Recording' when done."


def stop_recording():
    global recording, audio_frames
    recording = False
    time.sleep(0.5)  # Wait for thread to finish
    print("DEBUG: Stopped recording. Frames captured:", len(audio_frames))

    if not audio_frames:
        print("ERROR: No audio frames captured.")
        return "❌ No audio recorded. Please try again.", ""

    with tempfile.NamedTemporaryFile(delete=False, suffix=".wav") as temp_audio:
        wf = wave.open(temp_audio.name, 'wb')
        wf.setnchannels(CHANNELS)
        wf.setsampwidth(audio_interface.get_sample_size(FORMAT))
        wf.setframerate(RATE)
        wf.writeframes(b''.join(audio_frames))
        wf.close()

        try:
            with sr.AudioFile(temp_audio.name) as source:
                audio_data = recognizer.record(source)
                text = recognizer.recognize_google(audio_data)
                print("DEBUG: Transcription result:", text)
                return f"✅ Recording stopped. Transcribed: {text}", text
        except sr.UnknownValueError:
            print("ERROR: Could not understand audio.")
            return "❌ Could not understand audio. Please try again.", ""
        except sr.RequestError as e:
            print(f"ERROR: Speech recognition error: {e}")
            return f"❌ Speech recognition error: {e}", ""
        finally:
            import os
            os.unlink(temp_audio.name)


def text_to_speech(text):
    """Convert text to speech and return audio file"""
    try:
        # Use gTTS for better quality
        tts = gTTS(text=text, lang='en', slow=False)
        
        # Save to temporary file
        with tempfile.NamedTemporaryFile(delete=False, suffix=".mp3") as temp_file:
            tts.save(temp_file.name)
            return temp_file.name
    except Exception as e:
        print(f"Error in text-to-speech: {e}")
        return None


def chat(user_input, history):
    global document_context
    history = history or []
    user_input = normalize_input(user_input)

    insert_memory(user_input)

    # === Split into sub-questions ===
    sub_questions = re.split(r'\band\b|[,;]', user_input, flags=re.IGNORECASE)
    sub_questions = [q.strip() for q in sub_questions if q.strip()]

    answers = []

    for sq in sub_questions:
        # 1. Try document context if present
        if document_context:
            prompt = f"""### Instruction:
Answer ONLY using the information in the Document Context below. If the answer is not present, reply exactly: "The document does not contain this information."

### Document Context:
{document_context}

### User:
{sq}

### Response:"""
            inputs = tokenizer(prompt, return_tensors="pt").to(finetuned_model.device)
            with torch.no_grad():
                output = finetuned_model.generate(
                    **inputs,
                    max_new_tokens=128,
                    do_sample=False,
                    temperature=0.7,
                    top_p=0.9,
                    repetition_penalty=1.1,
                    pad_token_id=tokenizer.eos_token_id,
                    eos_token_id=tokenizer.eos_token_id
                )
            decoded = tokenizer.decode(output[0], skip_special_tokens=True)
            reply = decoded.split("### Response:")[-1].strip()
            if reply and "does not contain this information" not in reply:
                answers.append(reply)
                continue  # Skip to next sub-question if answered from document

        # 2. Otherwise, answer using facts/memory/database
        item = find_dataset_item(sq)
        is_new_fact = is_fact_statement(sq) and not (
        item and (
                item["instruction"].strip().lower() == sq.strip().lower() or
                item["instruction"].strip().lower() in sq.strip().lower() or
                sq.strip().lower() in item["instruction"].strip().lower()
        )
        )
        if is_new_fact:
            insert_relationship(sq)

        dataset_facts = []
        ds_item = find_dataset_item(sq)
        if ds_item:
            dataset_facts.append(ds_item["output"].strip())

        semantic_memories = search_similar_memory(sq, top_k=5)
        relationship_memories = search_relationship_memory(sq, top_k=5)

        def is_valid_fact(text):
            return "it department" not in text.lower()
        filtered_memories = [m for m in semantic_memories + relationship_memories if is_valid_fact(m)]

        memory_context = ""
        if dataset_facts:
            memory_context += "Dataset Facts:\n" + "\n".join(dataset_facts) + "\n"
        if filtered_memories:
            memory_context += "Related Memories:\n" + "\n".join(filtered_memories)

        acknowledge_instruction = ""
        if is_new_fact:
            acknowledge_instruction = (
                "The user just told you a new fact. Please start your response with a short, positive acknowledgment like "
                "'Thanks for letting me know!', 'That's awesome!', 'Got it!', or something similar, before continuing.\n\n"
            )

    prompt = f"""### Instruction:
{acknowledge_instruction}{sq}

### Facts:
{memory_context}

### User:
{sq}

### Response:"""
    inputs = tokenizer(prompt, return_tensors="pt").to(finetuned_model.device)
    with torch.no_grad():
        output = finetuned_model.generate(
            **inputs,
                max_new_tokens=128,
            do_sample=False,
            temperature=0.7,
            top_p=0.9,
            repetition_penalty=1.1,
            pad_token_id=tokenizer.eos_token_id,
            eos_token_id=tokenizer.eos_token_id
        )
    decoded = tokenizer.decode(output[0], skip_special_tokens=True)
    reply = decoded.split("### Response:")[-1].strip()
        answers.append(reply)

    # Combine all sub-answers
    final_reply = " ".join(answers)
    history.append({"role": "user", "content": user_input})
    history.append({"role": "assistant", "content": final_reply})

    # Store Q&A in document_memory if document was used
    if document_context:
        cursor.execute("INSERT INTO document_memory (question, answer) VALUES (%s, %s)", (user_input, final_reply))
        conn.commit()
        #document_context = None  # ✅ Reset document context for next upload

    return history, history


def extract_text_from_pdf(file):
    reader = PyPDF2.PdfReader(file)
    return "\n".join(page.extract_text() for page in reader.pages if page.extract_text())

def extract_text_from_docx(file):
    doc = docx.Document(file)
    return "\n".join(para.text for para in doc.paragraphs if para.text.strip())

def extract_text_from_doc(file):
    try:
        text = textract.process(file, extension='doc')
        return text.decode('utf-8')
    except Exception as e:
        return f"[Error reading DOC file: {e}]"

def extract_text_from_pptx(file):
    prs = pptx.Presentation(file)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def extract_text_from_excel(file, ext):
    # ext: .xls, .xlsx, .xlsm
    try:
        df = pd.read_excel(file, engine=None if ext == ".xls" else "openpyxl")
        return df.to_string(index=False)
    except Exception as e:
        return f"[Error reading Excel file: {e}]"

def extract_text_from_csv(file):
    try:
        df = pd.read_csv(file)
        return df.to_string(index=False)
    except Exception as e:
        return f"[Error reading CSV file: {e}]"

def extract_text_from_image(file):
    try:
        img = Image.open(file)
        text = pytesseract.image_to_string(img)
        return text
    except Exception as e:
        return f"[Error reading image: {e}]"

def extract_text_from_ppt(file):
    try:
        text = textract.process(file, extension='ppt')
        return text.decode('utf-8')
    except Exception as e:
        return f"[Error reading PPT file: {e}]"

def extract_text_from_file(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    with open(file_path, "rb") as f:
        if ext == ".pdf":
            return extract_text_from_pdf(f)
        elif ext == ".docx":
            return extract_text_from_docx(f)
        elif ext == ".txt":
            return f.read().decode("utf-8")
        # New: .doc (Word 97-2003)
        elif ext == ".doc":
            return extract_text_from_doc(file_path)
        # New: .pptx (PowerPoint)
        elif ext == ".pptx":
            return extract_text_from_pptx(file_path)
        # New: .ppt (old PowerPoint)
        elif ext == ".ppt":
            return extract_text_from_ppt(file_path)
        # New: .csv
        elif ext == ".csv":
            return extract_text_from_csv(file_path)
        # New: .xlsx, .xls, .xlsm
        elif ext in [".xlsx", ".xls", ".xlsm"]:
            return extract_text_from_excel(file_path, ext)
        # New: Images (jpg, jpeg, png)
        elif ext in [".jpg", ".jpeg", ".png"]:
            return extract_text_from_image(file_path)
    return ""


def clear_document_state():
    global document_context
    document_context = None
    return "🧹 Document cleared. Please upload a new one."


def handle_upload(file):
    global document_context
    document_context = extract_text_from_file(file.name)
    return "✅ Document uploaded and ready!"


def play_audio_response(text):
    """Generate and play audio for the AI response"""
    audio_file = text_to_speech(text)
    if audio_file:
        return audio_file
    else:
        return None


# 7. Gradio UI
with gr.Blocks() as demo:
    gr.Markdown("## 🤖 Chatbot with Structured + Semantic + Document Memory + Audio Features")
    
    with gr.Row():
        with gr.Column(scale=2):
            chatbot = gr.Chatbot(label="Chat", type="messages", height=400)
    msg = gr.Textbox(placeholder="Ask something...", label="Your Question")
    state = gr.State([])

            with gr.Row():
                submit_btn = gr.Button("Send", variant="primary")
                clear_btn = gr.Button("Clear Chat")
        
        with gr.Column(scale=1):
            gr.Markdown("### 🎤 Voice Input")
            with gr.Row():
                start_record_btn = gr.Button("🎤 Start Recording", variant="secondary")
                stop_record_btn = gr.Button("⏹️ Stop Recording", variant="stop")
            
            recording_status = gr.Textbox(label="Recording Status", interactive=False)
            transcribed_text = gr.Textbox(label="Transcribed Text", interactive=False)
            
            gr.Markdown("### 🔊 Audio Response")
            play_audio_btn = gr.Button("🔊 Play Last Response", variant="secondary")
            audio_output = gr.Audio(label="AI Response Audio", type="filepath")
            
            gr.Markdown("### 📄 Document Upload")
            file_upload = gr.File(label="📄 Upload Document (PDF, DOCX, TXT)")
            upload_output = gr.Textbox(label="Document Status", interactive=False)
    
    # Event handlers
    def process_audio_input():
        return "", ""  # Clear status and transcribed text
    
    def handle_chat_with_audio(user_input, history):
        if not user_input.strip():
            return history, history, None
        
        new_history, _ = chat(user_input, history)
        last_response = new_history[-1]["content"] if new_history else ""
        
        # Generate audio for the response
        audio_file = play_audio_response(last_response)
        
        return new_history, new_history, audio_file
    
    # Text input events
    submit_btn.click(
        handle_chat_with_audio, 
        inputs=[msg, state], 
        outputs=[chatbot, state, audio_output]
    ).then(
        lambda: "", inputs=None, outputs=msg
    ).then(
        clear_document_state, inputs=None, outputs=upload_output
    )
    
    msg.submit(
        handle_chat_with_audio, 
        inputs=[msg, state], 
        outputs=[chatbot, state, audio_output]
    ).then(
        lambda: "", inputs=None, outputs=msg
    ).then(
        clear_document_state, inputs=None, outputs=upload_output
    )
    
    # Audio recording events
    start_record_btn.click(
        start_recording, 
        inputs=None, 
        outputs=recording_status
    )
    
    stop_record_btn.click(
        stop_recording, 
        inputs=None, 
        outputs=[recording_status, transcribed_text]
    ).then(
        lambda: "", inputs=None, outputs=msg
    )

    # Play audio button
    play_audio_btn.click(
        lambda history: play_audio_response(history[-1]["content"]) if history else None,
        inputs=[state],
        outputs=audio_output
    )
    
    # File upload events
    file_upload.change(
        handle_upload, 
        inputs=file_upload, 
        outputs=upload_output
    )
    
    # Clear chat
    clear_btn.click(
        lambda: ([], []), 
        inputs=None, 
        outputs=[chatbot, state]
    ).then(
        process_audio_input, 
        inputs=None, 
        outputs=[recording_status, transcribed_text]
    )

demo.launch(server_name="0.0.0.0", server_port=8501, share=True)